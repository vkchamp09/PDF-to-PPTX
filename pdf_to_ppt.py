import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
import threading
import os
import sys

try:
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print(f"Missing required library: {e}")
    print("Please install required packages:")
    print("pip install PyMuPDF python-pptx")
    sys.exit(1)


class PDFtoPPTConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF to PowerPoint Converter")
        self.root.geometry("600x500")
        self.root.resizable(True, True)
        
        # Variables
        self.pdf_path = tk.StringVar()
        self.output_path = tk.StringVar()
        self.dpi_value = tk.IntVar(value=300)
        self.slide_width = tk.DoubleVar(value=10.0)
        self.is_converting = False
        
        self.setup_ui()
        
    def setup_ui(self):
        # Main frame
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # Configure grid weights
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        
        # Title
        title_label = ttk.Label(main_frame, text="PDF to PowerPoint Converter", 
                               font=("Arial", 16, "bold"))
        title_label.grid(row=0, column=0, columnspan=3, pady=(0, 20))
        
        # PDF File Selection
        ttk.Label(main_frame, text="Select PDF File:").grid(row=1, column=0, sticky=tk.W, pady=5)
        
        pdf_frame = ttk.Frame(main_frame)
        pdf_frame.grid(row=1, column=1, columnspan=2, sticky=(tk.W, tk.E), pady=5)
        pdf_frame.columnconfigure(0, weight=1)
        
        self.pdf_entry = ttk.Entry(pdf_frame, textvariable=self.pdf_path, width=50)
        self.pdf_entry.grid(row=0, column=0, sticky=(tk.W, tk.E), padx=(0, 5))
        
        ttk.Button(pdf_frame, text="Browse", command=self.browse_pdf).grid(row=0, column=1)
        
        # Output File Selection
        ttk.Label(main_frame, text="Output PPT File:").grid(row=2, column=0, sticky=tk.W, pady=5)
        
        output_frame = ttk.Frame(main_frame)
        output_frame.grid(row=2, column=1, columnspan=2, sticky=(tk.W, tk.E), pady=5)
        output_frame.columnconfigure(0, weight=1)
        
        self.output_entry = ttk.Entry(output_frame, textvariable=self.output_path, width=50)
        self.output_entry.grid(row=0, column=0, sticky=(tk.W, tk.E), padx=(0, 5))
        
        ttk.Button(output_frame, text="Browse", command=self.browse_output).grid(row=0, column=1)
        
        # Settings Frame
        settings_frame = ttk.LabelFrame(main_frame, text="Conversion Settings", padding="10")
        settings_frame.grid(row=3, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=10)
        settings_frame.columnconfigure(1, weight=1)
        
        # DPI Setting
        ttk.Label(settings_frame, text="DPI (Image Quality):").grid(row=0, column=0, sticky=tk.W, pady=5)
        dpi_frame = ttk.Frame(settings_frame)
        dpi_frame.grid(row=0, column=1, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Scale(dpi_frame, from_=150, to=600, variable=self.dpi_value, 
                 orient=tk.HORIZONTAL, length=200).grid(row=0, column=0, sticky=(tk.W, tk.E))
        
        self.dpi_label = ttk.Label(dpi_frame, text="300")
        self.dpi_label.grid(row=0, column=1, padx=(10, 0))
        
        # Update DPI label when scale changes
        self.dpi_value.trace('w', self.update_dpi_label)
        
        # Slide Width Setting
        ttk.Label(settings_frame, text="Slide Width (inches):").grid(row=1, column=0, sticky=tk.W, pady=5)
        width_frame = ttk.Frame(settings_frame)
        width_frame.grid(row=1, column=1, sticky=(tk.W, tk.E), pady=5)
        
        ttk.Scale(width_frame, from_=5.0, to=15.0, variable=self.slide_width, 
                 orient=tk.HORIZONTAL, length=200).grid(row=0, column=0, sticky=(tk.W, tk.E))
        
        self.width_label = ttk.Label(width_frame, text="10.0")
        self.width_label.grid(row=0, column=1, padx=(10, 0))
        
        # Update width label when scale changes
        self.slide_width.trace('w', self.update_width_label)
        
        # Convert Button
        self.convert_btn = ttk.Button(main_frame, text="Convert PDF to PPT", 
                                     command=self.start_conversion, style="Accent.TButton")
        self.convert_btn.grid(row=4, column=0, columnspan=3, pady=20)
        
        # Progress Bar
        self.progress = ttk.Progressbar(main_frame, mode='indeterminate')
        self.progress.grid(row=5, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        
        # Status Label
        self.status_label = ttk.Label(main_frame, text="Ready to convert", foreground="green")
        self.status_label.grid(row=6, column=0, columnspan=3, pady=5)
        
        # Log Text Area
        log_frame = ttk.LabelFrame(main_frame, text="Conversion Log", padding="5")
        log_frame.grid(row=7, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=10)
        log_frame.columnconfigure(0, weight=1)
        log_frame.rowconfigure(0, weight=1)
        
        self.log_text = ScrolledText(log_frame, height=8, wrap=tk.WORD)
        self.log_text.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # Configure main_frame row weights for resizing
        main_frame.rowconfigure(7, weight=1)
        
    def update_dpi_label(self, *args):
        self.dpi_label.config(text=str(self.dpi_value.get()))
        
    def update_width_label(self, *args):
        self.width_label.config(text=f"{self.slide_width.get():.1f}")
        
    def browse_pdf(self):
        filename = filedialog.askopenfilename(
            title="Select PDF File",
            filetypes=[("PDF files", "*.pdf"), ("All files", "*.*")]
        )
        if filename:
            self.pdf_path.set(filename)
            # Auto-generate output filename
            base_name = os.path.splitext(os.path.basename(filename))[0]
            output_dir = os.path.dirname(filename)
            suggested_output = os.path.join(output_dir, f"{base_name}.pptx")
            self.output_path.set(suggested_output)
    
    def browse_output(self):
        filename = filedialog.asksaveasfilename(
            title="Save PowerPoint As",
            defaultextension=".pptx",
            filetypes=[("PowerPoint files", "*.pptx"), ("All files", "*.*")]
        )
        if filename:
            self.output_path.set(filename)
    
    def log_message(self, message):
        """Add message to log with timestamp"""
        from datetime import datetime
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.insert(tk.END, f"[{timestamp}] {message}\n")
        self.log_text.see(tk.END)
        self.root.update_idletasks()
    
    def validate_inputs(self):
        """Validate user inputs before conversion"""
        if not self.pdf_path.get():
            messagebox.showerror("Error", "Please select a PDF file")
            return False
            
        if not os.path.exists(self.pdf_path.get()):
            messagebox.showerror("Error", "Selected PDF file does not exist")
            return False
            
        if not self.output_path.get():
            messagebox.showerror("Error", "Please specify output PowerPoint file")
            return False
            
        # Check if output directory exists
        output_dir = os.path.dirname(self.output_path.get())
        if output_dir and not os.path.exists(output_dir):
            messagebox.showerror("Error", f"Output directory does not exist: {output_dir}")
            return False
            
        return True
    
    def convert_pdf_to_ppt(self):
        """Main conversion function"""
        try:
            self.log_message("Starting PDF to PPT conversion...")
            
            # Open PDF
            self.log_message(f"Opening PDF: {self.pdf_path.get()}")
            doc = fitz.open(self.pdf_path.get())
            
            # Create PowerPoint presentation
            self.log_message("Creating PowerPoint presentation...")
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            
            # Calculate slide dimensions based on PDF aspect ratio
            first_page = doc[0]
            rect = first_page.rect
            pdf_width, pdf_height = rect.width, rect.height
            
            slide_width_in = self.slide_width.get()
            slide_height_in = slide_width_in * (pdf_height / pdf_width)
            
            prs.slide_width = Inches(slide_width_in)
            prs.slide_height = Inches(slide_height_in)
            
            self.log_message(f"Slide dimensions: {slide_width_in:.2f}\" x {slide_height_in:.2f}\"")
            self.log_message(f"Processing {len(doc)} pages at {self.dpi_value.get()} DPI...")
            
            # Process each page
            for page_num in range(len(doc)):
                if not self.is_converting:  # Check if conversion was cancelled
                    break
                    
                self.log_message(f"Processing page {page_num + 1}/{len(doc)}")
                
                page = doc[page_num]
                
                # Render PDF page to image
                pix = page.get_pixmap(dpi=self.dpi_value.get())
                img_path = f"temp_page_{page_num + 1}.png"
                pix.save(img_path)
                
                # Add new blank slide
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Insert the page image to fully fit the slide
                slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                       width=prs.slide_width, height=prs.slide_height)
                
                # Remove the temporary image
                os.remove(img_path)
                
                # Update progress
                self.root.update_idletasks()
            
            # Save PowerPoint file
            if self.is_converting:
                self.log_message(f"Saving PowerPoint file: {self.output_path.get()}")
                prs.save(self.output_path.get())
                
                doc.close()
                
                self.log_message("✅ Conversion completed successfully!")
                self.status_label.config(text="Conversion completed!", foreground="green")
                messagebox.showinfo("Success", f"PDF successfully converted to {self.output_path.get()}")
            else:
                doc.close()
                self.log_message("❌ Conversion cancelled by user")
                self.status_label.config(text="Conversion cancelled", foreground="orange")
                
        except Exception as e:
            self.log_message(f"❌ Error during conversion: {str(e)}")
            self.status_label.config(text="Conversion failed", foreground="red")
            messagebox.showerror("Error", f"Conversion failed: {str(e)}")
        
        finally:
            # Reset UI state
            self.is_converting = False
            self.progress.stop()
            self.convert_btn.config(text="Convert PDF to PPT", state="normal")
    
    def start_conversion(self):
        """Start the conversion process in a separate thread"""
        if self.is_converting:
            # Cancel conversion
            self.is_converting = False
            self.convert_btn.config(text="Cancelling...", state="disabled")
            return
            
        if not self.validate_inputs():
            return
            
        # Clear log
        self.log_text.delete(1.0, tk.END)
        
        # Update UI for conversion state
        self.is_converting = True
        self.convert_btn.config(text="Cancel Conversion")
        self.status_label.config(text="Converting...", foreground="blue")
        self.progress.start()
        
        # Start conversion in separate thread to keep UI responsive
        conversion_thread = threading.Thread(target=self.convert_pdf_to_ppt, daemon=True)
        conversion_thread.start()


class AboutDialog:
    def __init__(self, parent):
        self.dialog = tk.Toplevel(parent)
        self.dialog.title("About")
        self.dialog.geometry("400x300")
        self.dialog.resizable(False, False)
        self.dialog.transient(parent)
        self.dialog.grab_set()
        
        # Center the dialog
        self.dialog.geometry("+%d+%d" % (parent.winfo_rootx() + 50, parent.winfo_rooty() + 50))
        
        frame = ttk.Frame(self.dialog, padding="20")
        frame.pack(fill=tk.BOTH, expand=True)
        
        ttk.Label(frame, text="PDF to PowerPoint Converter", 
                 font=("Arial", 14, "bold")).pack(pady=(0, 10))
        
        info_text = """This tool converts PDF files to PowerPoint presentations.

Features:
• High-quality image conversion
• Maintains aspect ratio
• Customizable DPI settings
• Adjustable slide dimensions
• Real-time conversion progress

Requirements:
• PyMuPDF (fitz)
• python-pptx

Compatible with Python 3.8.16+"""
        
        ttk.Label(frame, text=info_text, justify=tk.LEFT).pack(pady=(0, 20))
        
        ttk.Button(frame, text="Close", command=self.dialog.destroy).pack()


def main():
    # Check Python version compatibility
    if sys.version_info < (3, 8):
        messagebox.showerror("Python Version Error", 
                           "This application requires Python 3.8 or higher")
        return
    
    root = tk.Tk()
    
    # Set the application icon (if available)
    try:
        # You can add an icon file here if you have one
        # root.iconbitmap("icon.ico")
        pass
    except:
        pass
    
    # Create menu bar
    menubar = tk.Menu(root)
    root.config(menu=menubar)
    
    # File menu
    file_menu = tk.Menu(menubar, tearoff=0)
    menubar.add_cascade(label="File", menu=file_menu)
    file_menu.add_separator()
    file_menu.add_command(label="Exit", command=root.quit)
    
    # Help menu
    help_menu = tk.Menu(menubar, tearoff=0)
    menubar.add_cascade(label="Help", menu=help_menu)
    help_menu.add_command(label="About", command=lambda: AboutDialog(root))
    
    # Create main application
    app = PDFtoPPTConverter(root)
    
    # Handle window closing
    def on_closing():
        if app.is_converting:
            if messagebox.askokcancel("Quit", "Conversion in progress. Do you want to quit?"):
                app.is_converting = False
                root.destroy()
        else:
            root.destroy()
    
    root.protocol("WM_DELETE_WINDOW", on_closing)
    
    # Start the GUI
    root.mainloop()


if __name__ == "__main__":
    main()
